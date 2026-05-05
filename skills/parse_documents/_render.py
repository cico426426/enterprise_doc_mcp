from pathlib import Path

import fitz
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def pdf_page_to_bytes(page: fitz.Page, scale: float = 2.0) -> bytes:
    """
    Render a PDF page as JPEG bytes.
    """
    try:
        matrix = fitz.Matrix(scale, scale)
        pixmap = page.get_pixmap(matrix=matrix, alpha=False)
        return pixmap.tobytes("jpeg")
    except Exception as exc:
        raise RuntimeError(f"Failed to render PDF page to JPEG bytes: {exc}") from exc


def pptx_slide_to_bytes(slide_index: int, pptx_path: Path) -> bytes | None:
    """
    Return the first embedded picture bytes in a slide.
    """
    presentation = Presentation(str(pptx_path))
    slide = presentation.slides[slide_index]
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return shape.image.blob
    return None
