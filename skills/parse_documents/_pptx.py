from pathlib import Path
from typing import Literal, TypedDict

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ._render import pptx_slide_to_bytes
from ._table import rows_to_markdown
from ._text import normalize_block
from ._vision import describe_image


PPTXRecord = TypedDict(
    "PPTXRecord",
    {
        "text": str,
        "source_file": str,
        "doc_type": Literal["pptx"],
        "kind": Literal["slide", "table"],
        "slide_number": int,
        "title": str,
        "has_table": bool,
        "has_visuals": bool,
        "vision_analyzed": bool,
    },
)


def _table_rows_from_shape(shape: object) -> list[list[str]]:
    table = getattr(shape, "table", None)
    if table is None:
        return []
    rows: list[list[str]] = []
    for row in table.rows:
        rows.append([cell.text or "" for cell in row.cells])
    return rows


def parse_pptx(
    path: Path,
    enable_vision: bool = True,
    vision_provider: str | None = None,
) -> list[PPTXRecord]:
    """
    Parse a PPTX into slide/table records.
    """
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise ValueError(f"Cannot open PPTX: {path}") from exc

    records: list[PPTXRecord] = []
    if len(prs.slides) == 0:
        return records

    for idx, slide in enumerate(prs.slides, start=1):
        title = normalize_block(slide.shapes.title.text) if slide.shapes.title and slide.shapes.title.text else ""
        text_chunks: list[str] = []
        has_table = False
        has_visual = False
        vision_analyzed = False

        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = normalize_block(shape.text_frame.text or "")
                if txt:
                    if title and txt == title:
                        continue
                    text_chunks.append(txt)

            if shape.has_table:
                rows = _table_rows_from_shape(shape)
                markdown = rows_to_markdown(rows)
                if markdown:
                    has_table = True
                    records.append(
                        PPTXRecord(
                            text=markdown,
                            source_file=path.name,
                            doc_type="pptx",
                            kind="table",
                            slide_number=idx,
                            title=title,
                            has_table=True,
                            has_visuals=False,
                            vision_analyzed=False,
                        )
                    )

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_visual = True

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = normalize_block(slide.notes_slide.notes_text_frame.text or "")
            if notes_text:
                text_chunks.append(f"[Notes]: {notes_text}")

        if enable_vision and has_visual:
            vision_analyzed = True
            image_bytes = pptx_slide_to_bytes(idx - 1, path)
            if image_bytes:
                result = describe_image(image_bytes, provider=vision_provider)
                if result:
                    text_chunks.append(
                        "\n".join(
                            [
                                "[Visual Content]",
                                f"Summary: {result.get('summary', '')}",
                                f"Charts: {result.get('charts', [])}",
                                f"Text: {result.get('text_content', '')}",
                            ]
                        )
                    )

        base_parts = []
        if title:
            base_parts.append(title)
        base_parts.extend(text_chunks)
        slide_text = "\n\n".join([p for p in base_parts if p]).strip()

        records.insert(
            len(records),
            PPTXRecord(
                text=slide_text,
                source_file=path.name,
                doc_type="pptx",
                kind="slide",
                slide_number=idx,
                title=title,
                has_table=has_table,
                has_visuals=has_visual,
                vision_analyzed=vision_analyzed,
            ),
        )

    return records
